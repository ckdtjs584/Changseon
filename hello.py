from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tkinter as tk
from tkinter import filedialog
import time
import tkinter.messagebox as msgbox





def extract_notes_from_pptx(file_path):
    # 프레젠테이션 열기
    presentation = Presentation(file_path)
    notes_data = []
    slide_texts = []

    # 각 슬라이드의 메모 추출
    for i, slide in enumerate(presentation.slides, start=1):
        notes_slide = slide.notes_slide
        slide_text = []
        # 각 슬라이드 안에 모든 쉐입을 돌면서 텍스트 프레임이 있는지 확인하여 리스트에 추가
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text.append(paragraph.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text.append(cell.text.strip())
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame and sub_shape.text.strip():
                        slide_text.append(sub_shape.text.strip())
        slide_texts.append((i, slide_text))

        # 각 슬라이드 안에 있는 메모를 추출
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text
        else:
            notes_text = "메모 없음"
        notes_data.append((i, notes_text))

    
    return (notes_data, slide_texts)


file_path = "D:/test.pptx"  # 대상 파워포인트 파일 경로
window = tk.Tk()

msgbox.showinfo("안내", "PPT 파일을 선택하세요.")
window.file = filedialog.askopenfile()
file_path = window.file.name
# file_path = r"D:\test.pptx" # PPT파일 경로


msgbox.showinfo("안내", f"내용이 작성될 txt파일을 선택하세요.\n(없다면 우클릭 > 새로만들기 > 텍스트 문서로 만들고 클릭)")
window.file = filedialog.askopenfile()
output_path = window.file.name
# output_path = r"D:\Output\output.txt" #
notes,slide_texts = extract_notes_from_pptx(file_path)



i = 0
with open(output_path, 'w', encoding='utf-8') as f:
    for slide_num, note in notes:
        
        f.write(f"(슬라이드 {slide_num})\n\n")
        f.write(f"메모  \n{note}\n\n")
        f.write(f"내부 텍스트  \n")
        j = 1

        for k in range(len(slide_texts[i][1])):
            f.write(f"({k+1})  {slide_texts[i][1][k]}\n")
            
        
        i += 1
        f.write(f"\n\n\n\n")
        f.write("-" * 40)
        f.write(f"\n\n\n\n")

exit()
